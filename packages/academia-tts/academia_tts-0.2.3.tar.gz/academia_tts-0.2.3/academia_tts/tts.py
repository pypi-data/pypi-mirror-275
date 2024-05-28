import os
import re
import torch
import datetime
from pptx import Presentation
from whisperspeech.pipeline import Pipeline
from pptx.util import Inches
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pydub import AudioSegment


class TTSProcessor:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.presentation = None
        self.notes = []

    def check_cuda(self):
        if not torch.cuda.is_available():
            raise BaseException("Currently `academia_tts` requires CUDA, make sure you are running this on a machine with GPU.")

    def extract_notes(self):
        self.presentation = Presentation(self.pptx_path)
        for idx, slide in enumerate(self.presentation.slides):
            notes_slide = slide.notes_slide
            if notes_slide:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    self.notes.append((idx + 1, notes_text, slide))

    def generate_filename(self, slide_number, notes_text):
        words = notes_text.split()
        first_six_words = "_".join(words[:6])
        safe_filename = re.sub(r'\W+', '', first_six_words)
        return f"{slide_number}_{safe_filename}.wav"

    def save_audio_files(self, output_dir, pipe):
        audio_files = []
        for slide_number, notes_text, _ in self.notes:
            sentences = notes_text.replace(":", ",")
            output_filename = self.generate_filename(slide_number, notes_text)
            output_path = os.path.join(output_dir, output_filename)
            pipe.generate_to_file(output_path, sentences)
            
            # Generate a silent audio segment with 0.4s duration
            silent_segment = AudioSegment.silent(duration=400)  # Duration in milliseconds
            
            # Load the original audio file
            original_audio = AudioSegment.from_file(output_path)
            
            # Concatenate the silent segment before and after the original audio
            final_audio = silent_segment + original_audio + silent_segment
            
            # Save the final audio to the output path
            final_audio.export(output_path, format="wav")
            
            # Append the final output path to the list of audio files
            audio_files.append(output_path)
            
        return audio_files

    def add_audio_to_pptx(self, audio_files):
        for (slide_number, _, slide), audio_file in zip(self.notes, audio_files):
            movie_shape = slide.shapes.add_movie(audio_file, Inches(0.5), Inches(0.5), Inches(1), Inches(1), poster_frame_image=None, mime_type='audio/wav')

            # Modify the XML to set the audio to play automatically
            movie_shape._element[0][0][0].set('vol', '300')  # Set volume
            movie_shape._element[0][0][0].set('mute', '0')  # Ensure it is not muted

    def process(self, output_dir):
        self.check_cuda()
        self.extract_notes()
        os.makedirs(output_dir, exist_ok=True)
        s2a_model = 'collabora/whisperspeech:s2a-q4-small-en+pl.model'
        pipe = Pipeline(s2a_ref=s2a_model)
        audio_files = self.save_audio_files(output_dir, pipe)
        self.add_audio_to_pptx(audio_files)
        return self.presentation